"""
Document text extraction and rebuilding for PDF, DOCX, PPTX.
"""
import os
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pathlib import Path


def _get_fonts_dir() -> str:
    if os.name == "nt":
        return os.path.join(os.environ.get("WINDIR", r"C:\Windows"), "Fonts")
    return "/usr/share/fonts/truetype"


def _find_font(name: str) -> str | None:
    """Find a system font file by name."""
    fonts_dir = _get_fonts_dir()
    path = os.path.join(fonts_dir, name)
    return path if os.path.exists(path) else None


def _select_font_for_text(text: str) -> dict:
    """Select the best font for the given text based on script detection.

    Returns a dict with either:
      {"fontname": "builtin-name"}           — use PyMuPDF built-in (CJK)
      {"fontfile": "path", "fontname": "x"}  — embed a system font
    """
    for ch in text:
        cp = ord(ch)
        # CJK Unified Ideographs / CJK Extension A / Bopomofo
        if 0x4E00 <= cp <= 0x9FFF or 0x3400 <= cp <= 0x4DBF:
            return {"fontname": "china-s"}
        # Hiragana / Katakana
        if 0x3040 <= cp <= 0x30FF or 0x31F0 <= cp <= 0x31FF:
            return {"fontname": "japan"}
        # Hangul
        if 0xAC00 <= cp <= 0xD7AF or 0x1100 <= cp <= 0x11FF:
            return {"fontname": "korea"}
        # Arabic / Hebrew
        if 0x0600 <= cp <= 0x06FF or 0x0590 <= cp <= 0x05FF:
            for name in ("tahoma.ttf", "segoeui.ttf", "arial.ttf"):
                f = _find_font(name)
                if f:
                    return {"fontfile": f, "fontname": "rtlf"}
        # Devanagari / Bengali / Tamil / other Indic
        if 0x0900 <= cp <= 0x0DFF:
            for name in ("nirmala.ttf", "mangal.ttf", "segoeui.ttf"):
                f = _find_font(name)
                if f:
                    return {"fontfile": f, "fontname": "indic"}
        # Thai / Lao
        if 0x0E00 <= cp <= 0x0EFF:
            for name in ("tahoma.ttf", "segoeui.ttf"):
                f = _find_font(name)
                if f:
                    return {"fontfile": f, "fontname": "thai"}
        # Cyrillic
        if 0x0400 <= cp <= 0x04FF:
            for name in ("arial.ttf", "segoeui.ttf"):
                f = _find_font(name)
                if f:
                    return {"fontfile": f, "fontname": "latf"}

    # Latin / fallback — prefer system TTF for proper rendering
    for name in ("arial.ttf", "segoeui.ttf"):
        f = _find_font(name)
        if f:
            return {"fontfile": f, "fontname": "latf"}

    # Linux fallbacks
    for path in (
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf",
    ):
        if os.path.exists(path):
            return {"fontfile": path, "fontname": "latf"}

    return {"fontname": "helv"}


def extract_text_from_pdf(file_path: str) -> list[dict]:
    """Extract text from a PDF at span level, preserving font metadata.

    Returns list of dicts per text block:
    {
        "page": int,
        "block_bbox": (x0, y0, x1, y1),
        "text": str,            # concatenated text from all spans (for translation)
        "spans": [              # individual spans with formatting info
            {
                "text": str,
                "bbox": (x0, y0, x1, y1),
                "font": str,
                "size": float,
                "color": int,
                "flags": int,   # 2=italic, 16=bold
            }
        ]
    }
    """
    doc = fitz.open(file_path)
    blocks = []

    for page_num, page in enumerate(doc):
        page_dict = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)

        for block in page_dict["blocks"]:
            if block["type"] != 0:  # skip image blocks
                continue

            block_spans = []
            block_text_parts = []

            for line in block["lines"]:
                line_texts = []
                for span in line["spans"]:
                    text = span["text"]
                    if not text.strip():
                        continue
                    block_spans.append({
                        "text": text,
                        "bbox": tuple(span["bbox"]),
                        "font": span["font"],
                        "size": span["size"],
                        "color": span["color"],
                        "flags": span["flags"],
                    })
                    line_texts.append(text.strip())
                if line_texts:
                    block_text_parts.append(" ".join(line_texts))

            if block_spans:
                blocks.append({
                    "page": page_num,
                    "block_bbox": tuple(block["bbox"]),
                    "text": " ".join(block_text_parts),
                    "spans": block_spans,
                })

    doc.close()
    return blocks


def extract_text_from_docx(file_path: str) -> list[dict]:
    """Extract paragraphs from a DOCX file.
    Returns list of {index, text, style} dicts.
    """
    doc = DocxDocument(file_path)
    blocks = []
    for i, para in enumerate(doc.paragraphs):
        if para.text.strip():
            blocks.append({
                "index": i,
                "text": para.text.strip(),
                "style": para.style.name if para.style else "Normal",
            })
    return blocks


def extract_text_from_pptx(file_path: str) -> list[dict]:
    """Extract text from a PPTX file.
    Returns list of {slide, shape_idx, text} dicts.
    """
    from pptx import Presentation
    prs = Presentation(file_path)
    blocks = []
    for slide_num, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                full_text = "\n".join(
                    para.text for para in shape.text_frame.paragraphs
                    if para.text.strip()
                )
                if full_text.strip():
                    blocks.append({
                        "slide": slide_num,
                        "shape_idx": shape_idx,
                        "text": full_text.strip(),
                    })
    return blocks


def _int_to_rgb(color_int: int) -> tuple:
    """Convert PyMuPDF integer color to (r, g, b) floats 0-1."""
    r = ((color_int >> 16) & 0xFF) / 255.0
    g = ((color_int >> 8) & 0xFF) / 255.0
    b = (color_int & 0xFF) / 255.0
    return (r, g, b)


# Map common PDF font families to system font variants: (regular, bold, italic, bold-italic)
_FONT_VARIANTS = {
    "arial": ("arial.ttf", "arialbd.ttf", "ariali.ttf", "arialbi.ttf"),
    "times": ("times.ttf", "timesbd.ttf", "timesi.ttf", "timesbi.ttf"),
    "segoe": ("segoeui.ttf", "segoeuib.ttf", "segoeuii.ttf", "segoeuiz.ttf"),
    "calibri": ("calibri.ttf", "calibrib.ttf", "calibrii.ttf", "calibriz.ttf"),
    "cambria": ("cambria.ttc", "cambriab.ttf", "cambriai.ttf", "cambriaz.ttf"),
    "verdana": ("verdana.ttf", "verdanab.ttf", "verdanai.ttf", "verdanaz.ttf"),
    "tahoma": ("tahoma.ttf", "tahomabd.ttf", "tahoma.ttf", "tahomabd.ttf"),
}


def _select_font_for_span(text: str, original_font: str, is_bold: bool, is_italic: bool) -> dict:
    """Select font for a span, trying to match the original font's bold/italic variant."""
    orig_lower = original_font.lower()
    for family, variants in _FONT_VARIANTS.items():
        if family in orig_lower:
            idx = (2 if is_italic else 0) + (1 if is_bold else 0)
            font_file = _find_font(variants[idx])
            if font_file:
                return {"fontfile": font_file, "fontname": f"f{family}{idx}"}
    # Fall back to script-based selection
    return _select_font_for_text(text)


def _split_at_word_boundary(text: str, target_chars: int) -> tuple:
    """Split text near target_chars at the nearest word boundary."""
    if target_chars >= len(text):
        return text, ""

    search_start = max(0, int(target_chars * 0.8))
    search_end = min(len(text), int(target_chars * 1.2))

    best_split = target_chars
    for i in range(search_start, search_end):
        if text[i] == " ":
            best_split = i
            break

    return text[:best_split].rstrip(), text[best_split:].lstrip()


def _insert_translated_spans(page, original_spans: list[dict], translated_text: str):
    """Distribute translated text across original span bboxes, preserving formatting."""
    total_original_chars = sum(len(s["text"]) for s in original_spans)
    if total_original_chars == 0:
        return

    remaining_text = translated_text
    for i, span in enumerate(original_spans):
        is_last = (i == len(original_spans) - 1)

        if is_last:
            span_text = remaining_text
        else:
            proportion = len(span["text"]) / total_original_chars
            char_count = max(1, round(proportion * len(translated_text)))
            span_text, remaining_text = _split_at_word_boundary(remaining_text, char_count)

        if not span_text.strip():
            continue

        bbox = fitz.Rect(span["bbox"])
        fontsize = span["size"]
        color = _int_to_rgb(span["color"])
        is_bold = bool(span["flags"] & 16)
        is_italic = bool(span["flags"] & 2)

        font_info = _select_font_for_span(span_text, span["font"], is_bold, is_italic)

        kwargs = dict(
            rect=bbox,
            buffer=span_text,
            fontsize=fontsize,
            align=fitz.TEXT_ALIGN_LEFT,
            color=color,
            **font_info,
        )

        rc = page.insert_textbox(**kwargs)

        # Text overflow: progressively reduce font size
        if rc < 0:
            for scale in (0.8, 0.65, 0.5):
                kwargs["fontsize"] = max(4, fontsize * scale)
                rc = page.insert_textbox(**kwargs)
                if rc >= 0:
                    break


def _insert_block_fallback(page, block: dict):
    """Fallback insertion for blocks without span info (backward compatibility)."""
    bbox = fitz.Rect(block.get("block_bbox", block.get("bbox", (0, 0, 100, 20))))
    text = block["text"]
    line_count = max(1, text.count("\n") + 1)
    fontsize = min(11, max(6, bbox.height / line_count * 0.7))

    font_info = _select_font_for_text(text)
    kwargs = dict(
        rect=bbox,
        buffer=text,
        fontsize=fontsize,
        align=fitz.TEXT_ALIGN_LEFT,
        color=(0, 0, 0),
        **font_info,
    )
    rc = page.insert_textbox(**kwargs)
    if rc < 0:
        kwargs["fontsize"] = max(5, fontsize * 0.65)
        page.insert_textbox(**kwargs)


def rebuild_pdf(
    original_path: str,
    translated_blocks: list[dict],
    output_path: str,
):
    """Create a translated PDF by redacting and reinserting at span level,
    preserving original font, size, color, and bold/italic formatting."""
    doc = fitz.open(original_path)

    # Group blocks by page
    pages_blocks: dict[int, list[dict]] = {}
    for block in translated_blocks:
        pages_blocks.setdefault(block["page"], []).append(block)

    for page_num, blocks in pages_blocks.items():
        page = doc[page_num]

        # Step 1: Redact all original text spans (collect all, then apply once per page)
        for block in blocks:
            spans = block.get("spans", [])
            if spans:
                for span in spans:
                    bbox = fitz.Rect(span["bbox"])
                    page.add_redact_annot(bbox, fill=(1, 1, 1))
            else:
                # Fallback: use block_bbox or bbox
                bbox = fitz.Rect(block.get("block_bbox", block.get("bbox", (0, 0, 100, 20))))
                page.add_redact_annot(bbox, fill=(1, 1, 1))
        page.apply_redactions()

        # Step 2: Insert translated text preserving formatting
        for block in blocks:
            spans = block.get("spans", [])
            translated_text = block["text"]

            if spans:
                _insert_translated_spans(page, spans, translated_text)
            else:
                _insert_block_fallback(page, block)

    doc.save(output_path)
    doc.close()


def rebuild_docx(
    original_path: str,
    translated_blocks: list[dict],
    output_path: str,
):
    """Create a translated DOCX by replacing paragraph text."""
    doc = DocxDocument(original_path)

    # Build index map
    block_map = {b["index"]: b["text"] for b in translated_blocks}

    for i, para in enumerate(doc.paragraphs):
        if i in block_map:
            # Preserve first run's formatting, replace text
            if para.runs:
                para.runs[0].text = block_map[i]
                for run in para.runs[1:]:
                    run.text = ""
            else:
                para.text = block_map[i]

    doc.save(output_path)


def rebuild_pptx(
    original_path: str,
    translated_blocks: list[dict],
    output_path: str,
):
    """Create a translated PPTX by replacing shape text."""
    from pptx import Presentation
    prs = Presentation(original_path)

    # Build lookup: (slide, shape_idx) -> translated text
    block_map = {(b["slide"], b["shape_idx"]): b["text"] for b in translated_blocks}

    for slide_num, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            key = (slide_num, shape_idx)
            if key in block_map and shape.has_text_frame:
                # Replace text while preserving first paragraph's formatting
                tf = shape.text_frame
                if tf.paragraphs:
                    first_para = tf.paragraphs[0]
                    if first_para.runs:
                        first_para.runs[0].text = block_map[key]
                        for run in first_para.runs[1:]:
                            run.text = ""
                    else:
                        first_para.text = block_map[key]
                    # Clear remaining paragraphs
                    for para in tf.paragraphs[1:]:
                        for run in para.runs:
                            run.text = ""

    prs.save(output_path)


def detect_doc_type(filename: str) -> str | None:
    """Detect document type from filename extension."""
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return "pdf"
    elif ext in (".docx", ".doc"):
        return "docx"
    elif ext in (".pptx", ".ppt"):
        return "pptx"
    return None
